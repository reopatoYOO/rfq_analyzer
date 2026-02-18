"""
PPT Parser - Extracts text and tables from PowerPoint files using python-pptx.
"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from models import ParsedDocument, ParsedPage

logger = logging.getLogger(__name__)


class PptParser:
    """Parses PowerPoint files and extracts text and tables per slide."""

    def parse(self, file_path: str) -> ParsedDocument:
        """
        Parse a PowerPoint file and return a ParsedDocument.

        Args:
            file_path: Path to the PPTX file.

        Returns:
            ParsedDocument with pages containing text and tables.
        """
        path = Path(file_path)
        doc = ParsedDocument(
            file_path=str(path.resolve()),
            file_name=path.name,
            file_type="pptx",
        )

        try:
            prs = Presentation(file_path)

            for i, slide in enumerate(prs.slides, start=1):
                text_parts = []
                tables = []

                for shape in slide.shapes:
                    # Extract text from text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                text_parts.append(para_text)

                    # Extract tables
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        tables.append(table_data)

                # Build full text
                full_text = "\n".join(text_parts)

                # Add table text representation
                if tables:
                    table_text_parts = []
                    for table in tables:
                        for row in table:
                            table_text_parts.append(" | ".join(row))
                    full_text += "\n\n[TABLE DATA]\n" + "\n".join(table_text_parts)

                parsed_page = ParsedPage(
                    page_number=i,
                    page_label=f"Slide {i}",
                    text=full_text,
                    tables=tables,
                )
                doc.pages.append(parsed_page)

            logger.info(f"Parsed PPT: {path.name} ({len(doc.pages)} slides)")

        except Exception as e:
            logger.error(f"Failed to parse PPT {path.name}: {e}")
            raise

        return doc
